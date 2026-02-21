#!/usr/bin/env python3
"""
PPT (PowerPoint) 生成功能測試
由 OpenClaw 建立 🐻
"""

import os
import sys
from datetime import datetime
from pathlib import Path

# 嘗試導入 PPT 套件
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx 未安裝")


def check_environment():
    """檢查環境"""
    print("🔧 檢查 PPT 生成環境...")
    
    if not PPTX_AVAILABLE:
        print("❌ python-pptx 套件未安裝")
        print("   請執行: pip install python-pptx")
        return False
    
    print("✅ python-pptx 套件可用")
    return True


def create_simple_ppt(output_path="測試簡報.pptx"):
    """建立簡單的 PPT 簡報"""
    if not PPTX_AVAILABLE:
        return False
    
    try:
        print(f"📊 建立 PPT 簡報: {output_path}")
        
        # 建立簡報物件
        prs = Presentation()
        
        # ===== 第1頁：標題頁 =====
        slide_layout = prs.slide_layouts[0]  # 標題投影片
        slide = prs.slides.add_slide(slide_layout)
        
        # 設定標題
        title = slide.shapes.title
        title.text = "美股盤前重點新聞彙整"
        
        # 設定副標題
        subtitle = slide.placeholders[1]
        subtitle.text = f"報告日期: {datetime.now().strftime('%Y年%m月%d日')}\n生成工具: OpenClaw 🐻"
        
        # ===== 第2頁：目錄 =====
        slide_layout = prs.slide_layouts[1]  # 標題與內容
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "簡報目錄"
        
        content = slide.placeholders[1]
        content.text = (
            "1. 市場概況\n"
            "2. 重點新聞分析\n"
            "3. 數據統計\n"
            "4. 風險提示\n"
            "5. 投資建議\n"
            "6. 總結"
        )
        
        # ===== 第3頁：市場概況 =====
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "市場概況"
        
        content = slide.placeholders[1]
        content.text = (
            "📈 主要指數表現:\n"
            "• 道瓊工業指數: 小幅上漲\n"
            "• 那斯達克指數: 科技股領漲\n"
            "• S&P 500指數: 穩健表現\n\n"
            "💰 資金流向:\n"
            "• 資金從成長股流向防禦型股票\n"
            "• 債市殖利率上升吸引部分資金\n"
            "• 投資策略趨向保守"
        )
        
        # ===== 第4頁：重點新聞分析 =====
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "重點新聞分析"
        
        content = slide.placeholders[1]
        content.text = (
            "🔥 高重要性新聞:\n"
            "• 美國1月就業報告超預期\n"
            "   - 非農數據大幅超出市場預期\n"
            "   - 可能推遲聯準會降息時程\n\n"
            "• 台積電ADR勁揚3-4%\n"
            "   - 帶動半導體類股走勢\n"
            "   - AI需求持續支撐股價\n\n"
            "📰 其他重要新聞:\n"
            "• AI投資熱潮持續\n"
            "• 企業財報季進行中\n"
            "• 資金流向防禦型股票"
        )
        
        # ===== 第5頁：數據統計（使用不同版面）=====
        slide_layout = prs.slide_layouts[5]  # 僅標題
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "數據統計"
        
        # 手動加入文字方塊
        from pptx.util import Inches
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # 第一段
        p = tf.add_paragraph()
        p.text = "📊 新聞統計:"
        p.font.bold = True
        
        # 第二段
        p = tf.add_paragraph()
        p.text = "• 總新聞數: 5則"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "• 高重要性: 2則"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 中重要性: 3則"
        p.level = 1
        
        # 第三段
        p = tf.add_paragraph()
        p.text = "⏰ 時間分布:"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = "• 最早發布: 08:45"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "• 最晚發布: 11:15"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• 平均間隔: 約30分鐘"
        p.level = 1
        
        # ===== 第6頁：風險提示 =====
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "風險提示"
        
        content = slide.placeholders[1]
        content.text = (
            "⚠️ 主要風險:\n"
            "• 強勁就業數據可能延後降息時程\n"
            "• AI相關股票估值偏高，需留意回調風險\n"
            "• 市場波動可能加大\n"
            "• 企業財報結果可能影響個股走勢\n"
            "• 資金流向變化可能導致板塊輪動\n\n"
            "🛡️ 風險管理建議:\n"
            "• 控制投資部位\n"
            "• 分散投資組合\n"
            "• 保持適當現金部位"
        )
        
        # ===== 第7頁：投資建議 =====
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "投資建議"
        
        content = slide.placeholders[1]
        content.text = (
            "🎯 短期建議:\n"
            "• 關注防禦型板塊（公用事業、必需消費品）\n"
            "• 密切關注後續經濟數據發布\n"
            "• 考慮價值型股票配置\n\n"
            "📈 中期建議:\n"
            "• 分散投資以降低風險\n"
            "• 關注AI相關供應鏈機會\n"
            "• 留意半導體產業復甦\n\n"
            "💰 長期建議:\n"
            "• 保持投資紀律\n"
            "• 定期檢視投資組合\n"
            "• 把握市場調整機會"
        )
        
        # ===== 第8頁：總結 =====
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "總結"
        
        content = slide.placeholders[1]
        content.text = (
            "✅ 重點回顧:\n"
            "1. 就業數據強勁，可能影響貨幣政策\n"
            "2. 台積電表現亮眼，帶動半導體類股\n"
            "3. AI投資熱潮持續，但需關注估值\n"
            "4. 資金流向趨向保守\n"
            "5. 風險與機會並存\n\n"
            "🔮 展望:\n"
            "• 密切關注聯準會政策動向\n"
            "• 留意企業財報表現\n"
            "• 把握結構性投資機會"
        )
        
        # ===== 第9頁：聯絡資訊 =====
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "感謝聆聽"
        
        subtitle = slide.placeholders[1]
        subtitle.text = (
            "報告生成時間:\n"
            f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n"
            "生成工具:\n"
            "OpenClaw 可愛助理 🐻\n\n"
            "注意事項:\n"
            "本報告僅供參考，不構成投資建議"
        )
        
        # 儲存簡報
        prs.save(output_path)
        
        file_size = os.path.getsize(output_path)
        print(f"✅ PPT 簡報建立完成: {file_size:,} bytes")
        print(f"   總頁數: {len(prs.slides)} 頁")
        
        return True
        
    except Exception as e:
        print(f"❌ PPT 建立失敗: {e}")
        import traceback
        traceback.print_exc()
        return False


def create_advanced_ppt(output_path="進階簡報.pptx"):
    """建立進階格式的 PPT（包含圖形、樣式）"""
    if not PPTX_AVAILABLE:
        return False
    
    try:
        print(f"🎨 建立進階 PPT 簡報: {output_path}")
        
        prs = Presentation()
        
        # ===== 自訂標題頁 =====
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # 標題
        title = slide.shapes.title
        title.text = "財經新聞分析報告"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        
        # 副標題
        subtitle = slide.placeholders[1]
        subtitle.text = "專業分析 · 數據驅動 · 智能生成"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x36, 0x60, 0x92)  # 藍色
        
        # ===== 圖表頁（示意）=====
        slide_layout = prs.slide_layouts[5]  # 僅標題
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "市場數據圖表"
        
        # 加入文字說明（實際應用可加入圖表）
        left = Inches(1)
        top = Inches(2)
        width = Inches(4)
        height = Inches(3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "📈 圖表示例位置\n\n此處可插入:\n• 股價走勢圖\n• 資金流向圖\n• 產業分布圖\n• 風險熱力圖"
        
        # 右側加入另一個文字方塊
        left = Inches(5)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        tf2.text = "💡 圖表生成說明\n\n使用工具:\n• matplotlib 生成圖表\n• Peekaboo 截取畫面\n• 插入為圖片\n• 自動更新數據"
        
        # ===== 時間軸頁 =====
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "新聞時間軸"
        
        # 建立時間軸文字
        timeline_text = (
            "🕗 08:45 - 台積電ADR勁揚\n"
            "   半導體類股走強\n\n"
            "🕘 09:00 - 就業數據發布\n"
            "   非農數據超預期\n\n"
            "🕤 09:45 - 資金流向分析\n"
            "   流向防禦型股票\n\n"
            "🕥 10:30 - AI投資趨勢\n"
            "   市場關注估值\n\n"
            "🕚 11:15 - 財報季更新\n"
            "   企業展望受關注"
        )
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.text = timeline_text
        
        # 儲存
        prs.save(output_path)
        
        file_size = os.path.getsize(output_path)
        print(f"✅ 進階 PPT 建立完成: {file_size:,} bytes")
        
        return True
        
    except Exception as e:
        print(f"❌ 進階 PPT 建立失敗: {e}")
        return False


def main():
    """主程式"""
    print("=" * 70)
    print("🐻 PPT (PowerPoint) 生成功能測試")
    print("=" * 70)
    
    # 檢查環境
    if not check_environment():
        return
    
    # 建立輸出目錄
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    
    # 生成簡單 PPT
    print("\n" + "=" * 70)
    simple_ppt = output_dir / "美股新聞簡報.pptx"
    simple_success = create_simple_ppt(simple_ppt)
    
    # 生成進階 PPT
    print("\n" + "=" * 70)
    advanced_ppt = output_dir / "進階財經簡報.pptx"
    advanced_success = create_advanced_ppt(advanced_ppt)
    
    # 結果總結
    print("\n" + "=" * 70)
    print("📋 PPT 生成結果總結:")
    
    if simple_success:
        size = os.path.getsize(simple_ppt)
        print(f"  ✅ 簡單簡報: {simple_ppt} ({size:,} bytes)")
    
    if advanced_success:
        size = os.path.getsize(advanced_ppt)
        print(f"  ✅ 進階簡報: {advanced_ppt} ({size:,} bytes)")
    
    if not simple_success and not advanced_success:
        print("  ❌ 所有 PPT 生成失敗")
    
    print("\n" + "=" * 70)
    print("💡 使用建議:")
    print("1. 可修改內容為實際新聞數據")
    print("2. 可加入圖表、圖片增強視覺效果")
    print("3. 可使用模板統一簡報風格")
    print("4. 可自動化定期生成簡報")
    print("\n🎉 PPT 生成功能測試完成！")
    print("=" * 70)


if __name__ == "__main__":
    main()