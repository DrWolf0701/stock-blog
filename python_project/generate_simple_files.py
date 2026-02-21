#!/usr/bin/env python3
"""
生成簡單、小流量的 PDF、PPT、Excel 檔案
由 OpenClaw 建立 🐻
"""

import os
import sys
from datetime import datetime
from pathlib import Path

# PDF 相關
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors

# Excel 相關
import pandas as pd

# PPT 相關
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def create_small_pdf(output_path):
    """建立小流量 PDF"""
    print(f"📄 建立小 PDF: {output_path}")
    
    try:
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        # 簡單標題
        c.setFont("Helvetica-Bold", 16)
        c.drawString(100, height - 100, "簡要財經快訊")
        
        c.setFont("Helvetica", 10)
        c.drawString(100, height - 130, f"{datetime.now().strftime('%Y-%m-%d')}")
        
        # 簡要內容
        y = height - 180
        items = [
            "📊 市場概況: 平穩",
            "💰 資金流向: 趨保守", 
            "⚠️ 風險等級: 中等",
            "🎯 建議: 謹慎觀察"
        ]
        
        for item in items:
            c.drawString(120, y, item)
            y -= 25
        
        # 頁尾
        c.setFont("Helvetica-Oblique", 8)
        c.drawString(100, 50, "OpenClaw 生成 🐻")
        
        c.save()
        
        size = os.path.getsize(output_path)
        print(f"✅ PDF 建立完成: {size:,} bytes")
        return True
        
    except Exception as e:
        print(f"❌ PDF 建立失敗: {e}")
        return False


def create_small_excel(output_path):
    """建立小流量 Excel"""
    print(f"📈 建立小 Excel: {output_path}")
    
    try:
        # 最小數據集
        data = {
            '項目': ['市場狀態', '風險等級', '建議動作', '更新時間'],
            '數值': ['平穩', '中等', '觀察', datetime.now().strftime('%H:%M')]
        }
        
        df = pd.DataFrame(data)
        
        # 使用最簡單的儲存方式
        df.to_excel(output_path, index=False, engine='openpyxl')
        
        size = os.path.getsize(output_path)
        print(f"✅ Excel 建立完成: {size:,} bytes")
        return True
        
    except Exception as e:
        print(f"❌ Excel 建立失敗: {e}")
        return False


def create_small_ppt(output_path):
    """建立小流量 PPT"""
    if not PPTX_AVAILABLE:
        print("❌ python-pptx 未安裝")
        return False
    
    print(f"📊 建立小 PPT: {output_path}")
    
    try:
        prs = Presentation()
        
        # 只有一頁簡報
        slide_layout = prs.slide_layouts[0]  # 標題頁
        slide = prs.slides.add_slide(slide_layout)
        
        # 標題
        title = slide.shapes.title
        title.text = "財經快訊"
        
        # 內容
        subtitle = slide.placeholders[1]
        subtitle.text = (
            f"更新: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"
            "要點:\n"
            "• 市場狀態: 平穩\n"
            "• 風險等級: 中等\n"
            "• 建議: 謹慎觀察\n\n"
            "OpenClaw 🐻"
        )
        
        prs.save(output_path)
        
        size = os.path.getsize(output_path)
        print(f"✅ PPT 建立完成: {size:,} bytes")
        return True
        
    except Exception as e:
        print(f"❌ PPT 建立失敗: {e}")
        return False


def create_minimal_files():
    """建立最小檔案集"""
    print("=" * 60)
    print("🐻 生成最小檔案集 (PDF, PPT, Excel)")
    print("=" * 60)
    
    # 建立輸出目錄
    output_dir = Path("minimal_output")
    output_dir.mkdir(exist_ok=True)
    
    # 檔案路徑
    pdf_path = output_dir / "minimal_report.pdf"
    excel_path = output_dir / "minimal_data.xlsx"
    ppt_path = output_dir / "minimal_slides.pptx"
    
    results = []
    
    # 生成 PDF
    print("\n1. PDF 生成:")
    pdf_success = create_small_pdf(pdf_path)
    results.append(("PDF", pdf_path, pdf_success))
    
    # 生成 Excel
    print("\n2. Excel 生成:")
    excel_success = create_small_excel(excel_path)
    results.append(("Excel", excel_path, excel_success))
    
    # 生成 PPT
    print("\n3. PPT 生成:")
    ppt_success = create_small_ppt(ppt_path)
    results.append(("PPT", ppt_path, ppt_success))
    
    # 結果總結
    print("\n" + "=" * 60)
    print("📋 生成結果:")
    
    total_size = 0
    for format_name, path, success in results:
        if success and path.exists():
            size = os.path.getsize(path)
            total_size += size
            print(f"  ✅ {format_name}: {path.name} ({size:,} bytes)")
        else:
            print(f"  ❌ {format_name}: 生成失敗")
    
    print(f"\n📊 總檔案大小: {total_size:,} bytes")
    print(f"   平均每個檔案: {total_size // len(results):,} bytes")
    
    print("\n" + "=" * 60)
    print("🎉 最小檔案集生成完成！")
    print("=" * 60)
    
    return all(success for _, _, success in results)


def main():
    """主程式"""
    success = create_minimal_files()
    
    if success:
        print("\n💡 檔案位置:")
        print("  cd ~/.openclaw/workspace/python_project/minimal_output")
        print("  ls -la")
        
        print("\n🚀 下一步:")
        print("  1. 可傳送這些檔案給你")
        print("  2. 可根據需求調整內容")
        print("  3. 可建立更完整的報告")
    else:
        print("\n⚠️ 部分檔案生成失敗，請檢查錯誤訊息")


if __name__ == "__main__":
    main()